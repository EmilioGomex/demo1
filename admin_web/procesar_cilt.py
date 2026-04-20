import os
import re
import glob
import json
from pathlib import Path

# ==========================================
# IMPORTANTE: Instalar dependencias antes de ejecutar
# pip install pymupdf python-pptx openpyxl fpdf2 pillow
# ==========================================
try:
    import fitz  # PyMuPDF
except ImportError:
    print("Falta pymupdf. Ejecuta: pip install pymupdf")
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    print("Falta python-pptx. Ejecuta: pip install python-pptx")
    Presentation = None

try:
    import openpyxl
except ImportError:
    print("Falta openpyxl. Ejecuta: pip install openpyxl")
    openpyxl = None

try:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos
except ImportError:
    print("Falta fpdf2. Ejecuta: pip install fpdf2")
    FPDF = None

from PIL import Image
import io

BASE_DIR = r"C:\eCILT\demo1\admin_web\CILT"
OUTPUT_DIR = r"C:\eCILT\demo1\admin_web\CILT_ESTANDARIZADO"

def ensure_dir(path):
    if not os.path.exists(path):
        os.makedirs(path)

# Reglas de estandarización
def estandarizar_duracion(texto):
    # Convierte horas a minutos
    if 'h' in texto.lower() or 'hora' in texto.lower():
        numeros = re.findall(r'\d+', texto)
        if numeros:
            return f"{int(numeros[0]) * 60} min"
    return texto

def estandarizar_frecuencia(texto):
    texto = texto.lower()
    if 'turno' in texto or 'diario' in texto: return 'Diario'
    if 'semana' in texto: return 'Semanal'
    if 'quincen' in texto: return 'Quincenal'
    if 'mes' in texto or 'mensual' in texto or 'necesidad' in texto: return 'Mensual'
    if 'trimestr' in texto: return 'Trimestral'
    if 'semestr' in texto: return 'Semestral'
    if 'año' in texto or 'ano' in texto: return 'Anual'
    return 'Mensual' # default

def determinar_estado_y_tipo(nombre_archivo, pasos_texto):
    texto_total = (nombre_archivo + " " + " ".join(pasos_texto)).lower()
    tipo = "Limpieza"
    estado = "Apagado"
    
    if 'lubricaci' in texto_total or 'lub' in texto_total:
        tipo = "Lubricación"
        estado = "Apagado"
    elif 'inspecci' in texto_total or 'insp' in texto_total:
        tipo = "Inspección"
        estado = "Encendido"
    elif 'reapriete' in texto_total or 'ajuste' in texto_total:
        tipo = "Reapriete"
        estado = "Apagado"
    
    return estado, tipo

def calcular_compartida(duracion_minutos):
    try:
        minutos = int(re.search(r'\d+', duracion_minutos).group())
        return "Si" if minutos < 15 else "No"
    except:
        return "Si"

def limpiar_pasos(pasos_raw):
    # Convertir a verbos infinitivos simples y limpiar
    pasos_limpios = []
    for paso in pasos_raw:
        if len(paso) < 5: continue
        # Quitar viñetas
        p = re.sub(r'^[\-\*\•\d\.\)]+\s*', '', paso).strip()
        if p and p.lower() not in [p.lower() for p in pasos_limpios]:
            pasos_limpios.append(p)
    return pasos_limpios

# Extracción de PDF
def extraer_pdf(filepath):
    pasos = []
    imagenes = []
    if not fitz: return pasos, imagenes
    doc = fitz.open(filepath)
    for i, page in enumerate(doc):
        text = page.get_text()
        if text: pasos.extend(text.split('\n'))
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            imagenes.append(image_bytes)
    return pasos, imagenes

# Extracción de PPTX
def extraer_pptx(filepath):
    pasos = []
    imagenes = []
    if not Presentation: return pasos, imagenes
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pasos.extend(shape.text.split('\n'))
                if hasattr(shape, "image"):
                    imagenes.append(shape.image.blob)
    except:
        pass
    return pasos, imagenes

# Extracción de XLSX
def extraer_xlsx(filepath):
    pasos = []
    imagenes = []
    if not openpyxl: return pasos, imagenes
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell and isinstance(cell, str):
                        pasos.extend(cell.split('\n'))
    except:
        pass
    return pasos, imagenes

class PDF(FPDF):
    def header(self):
        self.set_font("helvetica", "B", 15)
        self.cell(0, 10, "ESTÁNDAR CILT - OPERACIONES", border=False, new_x=XPos.LMARGIN, new_y=YPos.NEXT, align="C")
        self.ln(5)

def generar_pdf_maquina(maquina, tareas):
    if not FPDF: return
    pdf = PDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=12)
    
    for tarea in tareas:
        pdf.set_font("helvetica", "B", 14)
        pdf.set_text_color(0, 51, 153)
        pdf.cell(0, 10, f"1. {tarea['nombre_tarea']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_text_color(0, 0, 0)
        pdf.ln(2)
        
        pdf.set_font("helvetica", "B", 12)
        pdf.cell(0, 8, "2. Paso a Paso Operativo", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("helvetica", size=11)
        for i, paso in enumerate(tarea['pasos'], 1):
            # Encode para evitar errores charmap
            paso_text = paso.encode('latin-1', 'replace').decode('latin-1')
            # Evitar error FPDF de line wrapping cortando palabras gigantes (ej. guiones largos)
            paso_text = re.sub(r'\S{70,}', lambda m: m.group(0)[:67] + '...', paso_text)
            try:
                pdf.multi_cell(0, 6, f"{i}. {paso_text}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            except Exception:
                pass
        pdf.ln(5)
        
        pdf.set_font("helvetica", "B", 12)
        pdf.cell(0, 8, "3. Imágenes", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        if tarea['imagenes']:
            for idx, img_bytes in enumerate(tarea['imagenes'][:2]): # Limitar a 2 imagenes por tarea
                try:
                    img_path = os.path.join(OUTPUT_DIR, f"temp_{maquina}_{idx}.jpg")
                    image = Image.open(io.BytesIO(img_bytes))
                    if image.mode == 'RGBA': image = image.convert('RGB')
                    image.save(img_path)
                    pdf.image(img_path, w=100)
                    os.remove(img_path)
                except Exception as e:
                    pass
        else:
            pdf.set_font("helvetica", "I", 11)
            pdf.cell(0, 6, "Sin imagen", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(5)
        
        pdf.set_font("helvetica", "B", 12)
        pdf.cell(0, 8, "4. Detalles de la Tarea", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("helvetica", size=11)
        
        pdf.cell(0, 6, f"Nombre de la tarea: {tarea['nombre_tarea']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.cell(0, 6, f"Estado del equipo: {tarea['estado']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.cell(0, 6, f"Tipo: {tarea['tipo']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.cell(0, 6, f"Frecuencia: {tarea['frecuencia']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.cell(0, 6, f"Duración: {tarea['duracion']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.cell(0, 6, f"Tarea compartida: {tarea['compartida']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        
        pdf.ln(10)
        pdf.add_page()
    
    # Save the final machine PDF
    output_path = os.path.join(OUTPUT_DIR, f"CILT_{maquina}.pdf")
    try:
        pdf.output(output_path)
        print(f"Generado: {output_path}")
    except Exception as e:
        print(f"Error al guardar {output_path}: {e}")

def procesar_archivos():
    ensure_dir(OUTPUT_DIR)
    
    # Agrupar archivos por máquina
    archivos_por_maquina = {}
    for raiz, dirs, archivos in os.walk(BASE_DIR):
        for archivo in archivos:
            maquina = os.path.basename(raiz)
            ruta_completa = os.path.join(raiz, archivo)
            if maquina not in archivos_por_maquina:
                archivos_por_maquina[maquina] = []
            archivos_por_maquina[maquina].append(ruta_completa)
            
    # Procesar cada máquina
    for maquina, archivos in archivos_por_maquina.items():
        if maquina in ['Botellas', 'Latas']: continue
        print(f"Procesando máquina: {maquina} ({len(archivos)} archivos)")
        
        tareas_maquina = []
        for arch in archivos:
            ext = arch.lower().split('.')[-1]
            pasos_raw = []
            imagenes = []
            
            if ext == 'pdf':
                pasos_raw, imagenes = extraer_pdf(arch)
            elif ext == 'pptx':
                pasos_raw, imagenes = extraer_pptx(arch)
            elif ext in ['xlsx', 'xls']:
                pasos_raw, imagenes = extraer_xlsx(arch)
                
            if not pasos_raw:
                continue
                
            pasos_limpios = limpiar_pasos(pasos_raw)
            if not pasos_limpios: continue
            
            # Asignaciones según reglas
            estado, tipo = determinar_estado_y_tipo(os.path.basename(arch), pasos_limpios)
            duracion = "15 min" # Por defecto
            compartida = calcular_compartida(duracion)
            frecuencia = estandarizar_frecuencia("turno") # Por defecto diario/turno
            
            tarea = {
                "nombre_tarea": f"Procedimiento {tipo} - {maquina}",
                "pasos": pasos_limpios[:10], # Limitar a los primeros 10 pasos relevantes
                "imagenes": imagenes,
                "estado": estado,
                "tipo": tipo,
                "frecuencia": frecuencia,
                "duracion": duracion,
                "compartida": compartida
            }
            tareas_maquina.append(tarea)
            
        if tareas_maquina:
            generar_pdf_maquina(maquina, tareas_maquina)

if __name__ == "__main__":
    procesar_archivos()
